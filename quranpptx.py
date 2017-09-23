"""
quran.py

Usage:
  quran.py <sura> <outputfile> [--start=<ayat>] [--end=<ayat>] [--arabic-font=<arabicfont>]

Options:
  --start=<ayat>             Start ayat
  --end=<ayat>               End ayat
  --arabic-font=<arabicfont> Font to use for arabic [default: Calibri]
"""
import os
import docopt
import csv
from xml.dom import minidom
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE


SURA_NAMES = {
        1: "al-Faatiha (The Opening)",
        2: "al-Baqara (The Cow)",
        3: "aal-`Imraan (The Family of `Imraan)",
        4: "an-Nisaa' (The Women)",
        5: "al-Maaida (The Table)",
        6: "al-An`aam (The Cattle)",
        7: "al-A`araaf (The Heights)",
        8: "al-Anfaal (The Spoils of War)",
        9: "at-Tawba (The Repentance)",
        10: "Yunus (Jonas)",
        11: "Hud (Hud)",
        12: "Yusuf (Joseph)",
        13: "ar-Ra`ad (The Thunder)",
        14: "Ibrahim (Abraham)",
        15: "al-Hijr (The Rock)",
        16: "an-Nahl (The Bee)",
        17: "al-Israa' (The Night Journey)",
        18: "al-Kahf (The Cave)",
        19: "Maryam (Mary)",
        20: "Taa-haa (Taa-haa)",
        21: "al-Anbiyaa (The Prophets)",
        22: "al-Hajj (The Pilgrimage)",
        23: "al-Mu'minoon (The Believers)",
        24: "an-Noor (The Light)",
        25: "al-Furqaan (The Criterion)",
        26: "ash-Shu`araa (The Poets)",
        27: "an-Naml (The Ant)",
        28: "al-Qasas (The Stories)",
        29: "al-`Ankaboot (The Spider)",
        30: "ar-Room (The Romans)",
        31: "Luqman (Luqman)",
        32: "as-Sajda (The Prostration)",
        33: "al-Ahzaab (The Clans)",
        34: "Saba (Sheba)",
        35: "Faatir (The Originator)",
        36: "Yaseen (Yaseen)",
        37: "as-Saaffaat (Those drawn up in Ranks)",
        38: "Saad (The letter 'saad')",
        39: "az-Zumar (The Groups)",
        40: "al-Ghaafir (The Forgiver)",
        41: "Fussilat (Explained in detail)",
        42: "ash-Shura (Consultation)",
        43: "az-Zukhruf (Ornaments of gold)",
        44: "ad-Dukhaan (The Smoke)",
        45: "al-Jaathiya (Crouching)",
        46: "al-Ahqaf (The Dunes)",
        47: "Muhammad (Muhammad)",
        48: "al-Fath (The Victory)",
        49: "al-Hujraat (The Inner Apartments)",
        50: "Qaaf (The letter 'qaaf')",
        51: "adh-Dhaariyat (The Winnowing Winds)",
        52: "at-Tur (The Mount)",
        53: "an-Najm (The Star)",
        54: "al-Qamar (The Moon)",
        55: "ar-Rahman (The Beneficent)",
        56: "al-Waaqia (The Inevitable)",
        57: "al-Hadid (The Iron)",
        58: "al-Mujaadila (The Pleading Woman)",
        59: "al-Hashr (The Exile)",
        60: "al-Mumtahina (She that is to be examined)",
        61: "as-Saff (The Ranks)",
        62: "al-Jumu`a (Friday)",
        63: "al-Munaafiqoon (The Hypocrites)",
        64: "at-Taghaabun (Mutual Disillusion)",
        65: "at-Talaaq (Divorce)",
        66: "at-Tahrim (The Prohibition)",
        67: "al-Mulk (The Sovereignty)",
        68: "al-Qalam (The Pen)",
        69: "al-Haaqqa (The Reality)",
        70: "al-Ma`aarij (The Ascending Stairways)",
        71: "Nooh (Nooh)",
        72: "al-Jinn (the Jinn)",
        73: "al-Muzzammil (The Enshrouded One)",
        74: "al-Muddaththir (The Cloaked One)",
        75: "al-Qiyaama (The Resurrection)",
        76: "al-Insaan (Man)",
        77: "al-Mursalaat (The Emissaries)",
        78: "an-Naba (The Announcement)",
        79: "an-Naazi`aat (Those who drag forth)",
        80: "`Abasa (He frowned)",
        81: "at-Takwir (The Overthrowing)",
        82: "al-Infitaar (The Cleaving)",
        83: "al-Mutaffifin (Defrauding)",
        84: "al-Inshiqaaq (The Splitting Open)",
        85: "al-Burooj (The Constellations)",
        86: "at-Taariq (The Morning Star)",
        87: "al-A`ala (The Most High)",
        88: "al-Ghaashiya (The Overwhelming)",
        89: "al-Fajr (The Dawn)",
        90: "al-Balad (The City)",
        91: "ash-Shams (The Sun)",
        92: "al-Lail (The Night)",
        93: "ad-Dhuha (The Morning Hours)",
        94: "al-Inshira (The Consolation)",
        95: "at-Tin (The Fig)",
        96: "al-`Alaq (The Clot)",
        97: "al-Qadr (The Power, Fate)",
        98: "al-Bayyina (The Evidence)",
        99: "az-Zalzala (The Earthquake)",
        100: "al-`Aadiyaat (The Chargers)",
        101: "al-Qaari`a (The Calamity)",
        102: "at-Takaathur (Competition)",
        103: "al-`Asr (The Declining Day, Epoch)",
        104: "al-Humaza (The Traducer)",
        105: "al-Fil (The Elephant)",
        106: "Quraish (Quraysh)",
        107: "al-Maa`un (Almsgiving)",
        108: "al-Kawthar (Abundance)",
        109: "al-Kaafiroon (The Disbelievers)",
        110: "an-Nasr (Divine Support)",
        111: "al-Lahab (The Flame)",
        112: "al-Ikhlaas (Sincerity)",
        113: "al-Falaq (The Dawn)",
        114: "an-Naas (Mankind)" }

def load_suras(xmlfilename, translationfilename):
    the_dir = os.path.dirname(os.path.realpath(__file__))
    xmldoc = minidom.parse(os.path.join(the_dir, xmlfilename))
    suras = {}
    for sura in xmldoc.getElementsByTagName("sura"):
        this_sura = {}
        for ayat in sura.getElementsByTagName("aya"):
            this_sura[int(ayat.getAttribute("index"))] = dict(arabic=ayat.getAttribute("text"))
        suras[int(sura.getAttribute("index"))] = this_sura
    for row in csv.reader(open(os.path.join(the_dir, translationfilename), "r")):
        sura = int(row[1])
        ayat = int(row[2])
        english = row[3]
        suras[sura][ayat]["english"] = english
    return suras

def create_presentation(sura_number, outputfile, start=None, end=None, arabic_font="Calibri"):
    suras = load_suras("quran-uthmani.xml", "shakir_table.csv")
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Sura %s" % (SURA_NAMES[sura_number],)
    if start:
        subtitle.txt = "Ayat %s to %s" % (start, end)
    blank_slide_layout = prs.slide_layouts[6]
    if sura_number != 9 and not (sura_number == 1 and (start is None or int(start) == 1)):
        # add bismillah
        ayat = suras[1][1]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text_frame.paragraphs[0].text
        subtitle = slide.placeholders[1]

        title.text = ayat["arabic"]
        title.text_frame.paragraphs[0].font.name = 'Arial'
        subtitle.text = ayat["english"]

    for number, ayat in suras[sura_number].iteritems():
        if start is None or (number >= int(start) and number <= int(end)):
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = ayat["arabic"]
            title.text_frame.paragraphs[0].font.name = 'Arial'
            subtitle.text = ayat["english"] + " [{}]".format(number)

    prs.save(outputfile)

if __name__ == '__main__':
    arguments = docopt.docopt(__doc__, version='quran.py 0.1')
    create_presentation(int(arguments["<sura>"]), arguments["<outputfile>"], arguments["--start"], arguments["--end"], arguments["--arabic-font"])
